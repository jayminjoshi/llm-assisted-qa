from .base_extractor import BaseExtractor
from typing import Dict, List, Union, Any
from pptx import Presentation

import uuid
import os

class PPTExtractor(BaseExtractor):
    def __init__(self):
        """Initialize the PowerPoint extractor"""
        super().__init__()
        self.supported_formats = ['.pptx', '.ppt']

    def validate_file(self, file_path: str) -> bool:
        """
        Validate if the file is a PowerPoint presentation
        :param file_path: Path to the PowerPoint file
        :return: Boolean indicating if file is valid
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
            
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format. Supported formats: {self.supported_formats}")
        
        return True

    def extract_text_from_shape(self, shape) -> str:
        """
        Extract text from a shape object with debugging
        :param shape: PowerPoint shape object
        :return: Extracted text
        """
        text = ""
        
        try:
            # Print shape type for debugging
            shape_type = shape.shape_type
            
            # Handle different shape types
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        row_texts.append(cell.text.strip())
                    text += " | ".join(row_texts) + "\n"
            
            elif hasattr(shape, "shapes"):  # Group shape
                for sub_shape in shape.shapes:
                    text += self.extract_text_from_shape(sub_shape) + "\n"
                    
            # Handle other specific shape types if needed
            elif shape_type == 7:  # Media
                text = "[Media Content]"
            elif shape_type == 13:  # Picture
                text = "[Picture]"
                
        except Exception as e:
            print(f"Warning: Could not extract text from shape type {getattr(shape, 'shape_type', 'unknown')}: {str(e)}")
            
        return text.strip()

    def extract_content(self, file_path: str) -> List[Dict[str, Union[int, str]]]:
        """
        Extract content from PowerPoint presentation
        :param file_path: Path to the PowerPoint file
        :return: List of dictionaries containing slide numbers and content
        """
        self.validate_file(file_path)
        presentation = Presentation(file_path)
        slides_content = []

        for slide_number, slide in enumerate(presentation.slides, 1):
            slide_content = {
                'chunk_number': slide_number,
                'content': [],
                'notes': '',
                'total_chunks': len(presentation.slides)
            }

            # Extract content from shapes
            for shape in slide.shapes:
                text = self.extract_text_from_shape(shape)
                if text:
                    slide_content['content'].append(text)

            # Extract notes if they exist
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for shape in notes_slide.shapes:
                    text = self.extract_text_from_shape(shape)
                    if text:
                        slide_content['notes'] += text + "\n"

            # Clean up the content
            slide_content['content'] = '\n'.join(slide_content['content'])
            slide_content['notes'] = slide_content['notes'].strip()
            
            slides_content.append(slide_content)

        return slides_content
    
    def extract_documents(
        self, 
        file_path: str, 
        project_id: int, 
        user_id: int,
        file_id: int
    ) -> List[Dict[str, Any]]:
        """
        Extract content from PowerPoint presentation in Qdrant format
        
        Args:
            file_path (str): Path to the PowerPoint file
            project_id (int): Project identifier
            user_id (int): User identifier
            file_id (int): File identifier
            
        Returns:
            List[Dict]: List of documents with page_content and metadata
        """
        slides_content = self.extract_content(file_path)
        documents = []
        
        for slide in slides_content:
            # Combine content and notes if notes exist
            page_content = slide['content']
            if slide['notes']:
                page_content += f"\nNotes: \n{slide['notes']}"
            
            # Create document with content and metadata
            document = {
                "vector_id": str(uuid.uuid4()),
                "page_content": page_content,
                "metadata": {
                    "project_id": project_id,
                    "user_id": user_id,
                    "file_id": file_id,
                    "file_type": "pptx",
                    "chunk_number": slide['chunk_number'],
                    "total_chunks": slide['total_chunks']
                }
            }
            
            self.insert_vector(document)
            documents.append(document)
        
        return documents